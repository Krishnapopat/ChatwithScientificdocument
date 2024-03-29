from flask import Flask, request, jsonify, render_template
from flask_cors import CORS
import fitz  # PyMuPDF
import os
import requests
from dotenv import load_dotenv
from sentence_transformers import SentenceTransformer, util
from pptx import Presentation
import docx
# import textract
import re
import cv2
import easyocr
import matplotlib.pyplot as plt
from IPython.display import Image
from pylab import rcParams


load_dotenv()

app = Flask(__name__)
CORS(app)

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
sentence_model = SentenceTransformer('all-MiniLM-L6-v2')
document_text = ""

def extract_text_from_pdf(pdf_content):
    with fitz.open(stream=pdf_content, filetype="pdf") as doc:
        text = ""
        for page in doc:
            text += page.get_text()
    return text


def extract_text_and_images_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    reader = easyocr.Reader(['en'])

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
            if hasattr(shape, 'image'):  
                image = shape.image
                image_bytes = image.blob
                output = reader.readtext(image_bytes)
                for detection in output:
                    text += detection[1] + "\n"

    return text

def extract_text_from_txt(txt_file):
    """
    Extract text from a .txt file.
    
    Args:
    - txt_file: File object of the .txt file.
    
    Returns:
    - text: Extracted text from the .txt file
    """
    try:
        # Assuming txt_file is a file object, read its content directly
        text = txt_file.read().decode('utf-8')
        return text
    except Exception as e:
        print(f"Error extracting text from .txt file: {e}")
        return None
def extract_text_from_docx(docx_file):
    """
    Extract text from a .docx file.
    
    Args:
    - docx_file: Path to the .docx file
    
    Returns:
    - text: Extracted text from the .docx file
    """
    try:
        doc = docx.Document(docx_file)
        text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        return text
    except Exception as e:
        print(f"Error extracting text from .docx file: {e}")
        return None
    
def extract_text_from_tex(tex_content):
    try:
        # Decode the content if it's bytes
        if isinstance(tex_content, bytes):
            tex_content = tex_content.decode('utf-8')
        
        # Remove comments
        tex_content = re.sub(r'%.*?\n', '', tex_content)
        
        # Replace common LaTeX symbols with their text equivalent
        # This might need to be extended based on specific needs
        replacements = {
            '\\&': '&',
            '\\%': '%',
            '\\_': '_',
            '\\#': '#',
            '\\{': '{',
            '\\}': '}',
            '\\textbf{': '',
            '\\textit{': '',
            '\\usepackage{': '',
            '\\begin{': '',
            '\\end{': '',
            # Add more replacements as needed
        }
        for key, value in replacements.items():
            tex_content = tex_content.replace(key, value)
        
        # Remove LaTeX commands with optional arguments
        tex_content = re.sub(r'\\[a-zA-Z]+\[[^\]]*\]', '', tex_content)
        
        # Remove LaTeX commands with mandatory arguments
        tex_content = re.sub(r'\\[a-zA-Z]+(\{[^}]*\})+', '', tex_content)
        
        # Remove remaining LaTeX commands without arguments
        tex_content = re.sub(r'\\[a-zA-Z]+', '', tex_content)
        
        # Remove braces around single words or spaces (from removed commands)
        tex_content = re.sub(r'\{([^}]*)\}', '\\1', tex_content)
        
        # Remove extra spaces and newlines
        tex_content = re.sub(r'\s+', ' ', tex_content).strip()
        
        return tex_content
    except Exception as e:
        print(f"Error extracting text from .tex file: {e}")
        return None


# def extract_text_from_doc(doc_file):
#     """
#     Extract text from a .doc file.
    
#     Args:
#     - doc_file: Path to the .doc file
    
#     Returns:
#     - text: Extracted text from the .doc file
#     """
#     try:
#         # Use textract to extract text from .doc files
#         text = textract.process(doc_file).decode("utf-8")
#         return text
#     except Exception as e:
#         print(f"Error extracting text from .doc file: {e}")
#         return None


@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'document' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['document']    
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    if file:
        file_extension = file.filename.rsplit('.', 1)[1].lower()
        
        global document_text
        if file_extension == 'pdf':
            document_text = extract_text_from_pdf(file.read())
            # print(document_text)
            return jsonify({"text": document_text}), 200
        elif file_extension == 'docx':
            document_text = extract_text_from_docx(file)
            # print(document_text)
            return jsonify({"text": document_text}), 200
        elif file_extension == 'tex':
            document_text = extract_text_from_tex(file.read())
            print(document_text)
            return jsonify({"text": document_text}), 200
        # elif file_extension == 'doc':
        #     document_text = extract_text_from_doc(file)
        #     # print(document_text)
        #     return jsonify({"text": document_text}), 200
        elif file_extension == 'txt':
            document_text = extract_text_from_txt(file)
            # print(document_text)
            return jsonify({"text": document_text}), 200
        elif file_extension == 'pptx':
            document_text = extract_text_and_images_from_pptx(file) 
            # print(document_text) 
            return jsonify({"text": document_text}), 200
        elif file_extension == 'ppt':
            return jsonify({"error": "Unsupported file type '.ppt'. Please convert to '.pptx' before uploading."}), 400
        else:
            return jsonify({"error": "Unsupported file type"}), 400
    else:
        return jsonify({"error": "Unsupported file type"}), 400


@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.get_json()
    question =  data.get("question") + "Provide precise answer in context of above document also shortly state apporopriate pages or title of the page from the document from where the query is answered if the question is related to the document."
    
    global document_text
    # Check if document text is available
    if not document_text:
        return jsonify({"answer": "Please upload a document first."}), 400
    
    # Use ChatGPT API to generate response
    response = requests.post(
        "https://api.openai.com/v1/chat/completions",
        headers={
            "Authorization": f"Bearer {OPENAI_API_KEY}",
            "Content-Type": "application/json",
        },
        json={
            "model": "gpt-3.5-turbo",
            "messages": [
                {"role": "user", "content": document_text},
                {"role": "assistant", "content": question},
            ],
        }
    )
    print(response.text)
    if response.status_code == 200:
        # answer = response.json().get("choices")[0]['message']['content'].get("text", "No answer found.").strip()
        answer = response.json().get("choices")[0]['message']['content']
    else:
        answer = "Failed to get a response from OpenAI API."

    return jsonify({"answer": answer}), 200

if __name__ == '__main__':
    port = int(os.getenv("PORT", 5000))  # Default to 5000 if PORT is not specified
    app.run(host='0.0.0.0', port=port, debug=True)
